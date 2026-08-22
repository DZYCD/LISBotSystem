import os, zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document


def _extract_text_from_slide(slide_xml):
    root = ET.fromstring(slide_xml)
    texts = []
    for t in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
        if t.text:
            texts.append(t.text)
    return '\n'.join(texts) if texts else ''


def _extract_notes_from_slide(notes_xml):
    root = ET.fromstring(notes_xml)
    texts = []
    for t in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
        if t.text:
            texts.append(t.text)
    return '\n'.join(texts) if texts else ''


def _parse_pptx(filepath):
    if not os.path.exists(filepath):
        return {"error": f"文件不存在: {filepath}"}
    if not filepath.lower().endswith('.pptx'):
        return {"error": f"不是 PPTX 文件: {filepath}"}
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            slide_files = sorted([f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])
            if not slide_files:
                return {"error": "无法找到幻灯片内容"}
            result = {"filename": os.path.basename(filepath), "slide_count": len(slide_files), "slides": []}
            for i, slide_file in enumerate(slide_files, 1):
                text = _extract_text_from_slide(z.read(slide_file))
                notes = ''
                notes_file = slide_file.replace('slides/slide', 'notesSlides/notesSlide')
                try:
                    notes = _extract_notes_from_slide(z.read(notes_file))
                except KeyError:
                    pass
                result["slides"].append({"index": i, "text": text, "notes": notes})
            return result
    except zipfile.BadZipFile:
        return {"error": f"文件损坏或不是有效的 PPTX: {filepath}"}
    except Exception as e:
        return {"error": f"解析失败: {str(e)}"}


def _format_slides(data):
    if "error" in data:
        return f"❌ {data['error']}"
    lines = [f"📊 文件: {data['filename']}", f"📄 共 {data['slide_count']} 张幻灯片", '=' * 50]
    for slide in data["slides"]:
        lines.append(f"\n【幻灯片 {slide['index']}】")
        lines.append(slide["text"] if slide["text"] else "[空幻灯片或无文本内容]")
        if slide["notes"]:
            lines.append(f"\n📝 备注:\n{slide['notes']}")
    return '\n'.join(lines)


def extract_ppt(path):
    return _format_slides(_parse_pptx(path))


def extract_batch(directory, recursive=False):
    if not os.path.isdir(directory):
        return f"错误：目录不存在: {directory}"
    if recursive:
        pptx_files = []
        for root, dirs, files in os.walk(directory):
            for f in files:
                if f.lower().endswith('.pptx') and not f.startswith('~$'):
                    pptx_files.append(os.path.join(root, f))
    else:
        pptx_files = [os.path.join(directory, f) for f in os.listdir(directory) if f.lower().endswith('.pptx') and not f.startswith('~$')]
    pptx_files.sort()
    if not pptx_files:
        return f"在目录中未找到 PPTX 文件: {directory}"
    results = []
    success = 0
    failed = 0
    for pptx_path in pptx_files:
        data = _parse_pptx(pptx_path)
        if "error" in data:
            results.append(f"❌ {os.path.basename(pptx_path)}: {data['error']}")
            failed += 1
        else:
            results.append(f"✅ {data['filename']} ({data['slide_count']} 页)")
            success += 1
    return f"批量提取完成：成功 {success} 个，失败 {failed} 个，共 {len(pptx_files)} 个文件\n" + '\n'.join(results)


def _extract_paragraphs(filepath):
    prs = Presentation(filepath)
    result = []
    for si, slide in enumerate(prs.slides, 1):
        for shi, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if text:
                        result.append({"slide": si, "shape": shi, "para": pi, "text": text})
    return result


def _set_paragraph_text(paragraph, text):
    for run in paragraph.runs:
        run.text = ''
    if paragraph.runs:
        paragraph.runs[0].text = text
    else:
        paragraph.add_run().text = text


def edit_ppt(path, action, old_text='', new_text='', case_sensitive=False, replacements=None, replace_mode='run', title_text='', body_text=''):
    if not os.path.exists(path):
        return f"错误：文件不存在: {path}"
    if not path.lower().endswith('.pptx'):
        return f"错误：不是 PPTX 文件: {path}"
    try:
        prs = Presentation(path)
        changed = False

        if action == "list_paragraphs":
            paras = _extract_paragraphs(path)
            lines = [f"幻灯片{p['slide']} 形状{p['shape']} 段落{p['para']}: {p['text']}" for p in paras]
            return "共有 " + str(len(paras)) + " 个非空段落：\n" + "\n".join(lines)

        if action not in ("replace_text", "add_slide", "list_paragraphs"):
            return "错误：action 必须为 replace_text, add_slide 或 list_paragraphs"

        if action == "replace_text":
            if replacements is None:
                replacements = [{"old_text": old_text, "new_text": new_text, "case_sensitive": case_sensitive}]
            for r in replacements:
                o_text = r.get("old_text", "")
                n_text = r.get("new_text", "")
                cs = r.get("case_sensitive", False)
                if not o_text:
                    continue
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            if replace_mode == "paragraph":
                                full_text = paragraph.text
                                if cs:
                                    if o_text in full_text:
                                        _set_paragraph_text(paragraph, full_text.replace(o_text, n_text))
                                        changed = True
                                else:
                                    if o_text.lower() in full_text.lower():
                                        _set_paragraph_text(paragraph, full_text.replace(o_text, n_text))
                                        changed = True
                            else:
                                for run in paragraph.runs:
                                    if cs:
                                        if o_text in run.text:
                                            run.text = run.text.replace(o_text, n_text)
                                            changed = True
                                    else:
                                        if o_text.lower() in run.text.lower():
                                            run.text = run.text.replace(o_text, n_text)
                                            changed = True
            if not changed:
                return "未找到匹配文本，未做修改"
        elif action == "add_slide":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            if body_text:
                try:
                    body_shape = slide.placeholders[1]
                    body_shape.text = body_text
                except (IndexError, KeyError):
                    pass
            changed = True
        prs.save(path)
        return f"✅ 修改成功：{action} 操作已完成"
    except Exception as e:
        return f"❌ 修改失败：{str(e)}"


def _read_docx_paragraphs(filepath):
    doc = Document(filepath)
    result = []
    for i, para in enumerate(doc.paragraphs, 1):
        text = para.text.strip()
        if text:
            result.append({"index": i, "text": text})
    return result


def _read_docx_runs(filepath):
    doc = Document(filepath)
    result = []
    for i, para in enumerate(doc.paragraphs, 1):
        text = para.text.strip()
        if not text:
            continue
        runs_info = []
        for j, run in enumerate(para.runs):
            if not run.text.strip():
                continue
            runs_info.append({"run_index": j, "text": run.text, "bold": run.bold if run.bold is not None else False, "italic": run.italic if run.italic is not None else False})
        if runs_info:
            result.append({"paragraph_index": i, "runs": runs_info})
    return result


def read_docx(path, detail_level='paragraph'):
    if not os.path.exists(path):
        return f"错误：文件不存在: {path}"
    if not path.lower().endswith('.docx'):
        return f"错误：不是 DOCX 文件: {path}"
    try:
        if detail_level == "run":
            runs_data = _read_docx_runs(path)
            lines = []
            for r in runs_data:
                lines.append(f"段落{r['paragraph_index']}:")
                for run in r['runs']:
                    flags = []
                    if run['bold']: flags.append("粗体")
                    if run['italic']: flags.append("斜体")
                    flag_str = f"[{','.join(flags)}]" if flags else ""
                    lines.append(f"   Run{run['run_index']} {flag_str}: {run['text']}")
            return '\n'.join(lines) if lines else "无内容"
        else:
            paras = _read_docx_paragraphs(path)
            lines = [f"段落{p['index']}: {p['text']}" for p in paras]
            return f"共 {len(paras)} 个非空段落：\n" + "\n".join(lines)
    except Exception as e:
        return f"❌ 读取失败：{str(e)}"


def edit_docx(path, action, old_text='', new_text='', case_sensitive=False, text=''):
    if not os.path.exists(path):
        return f"错误：文件不存在: {path}"
    if not path.lower().endswith('.docx'):
        return f"错误：不是 DOCX 文件: {path}"
    if action not in ("replace_text", "add_paragraph"):
        return "错误：action 必须为 replace_text 或 add_paragraph"
    try:
        doc = Document(path)
        changed = False
        if action == "replace_text":
            if not old_text:
                return "错误：replace_text 需要提供 old_text"
            for para in doc.paragraphs:
                if case_sensitive:
                    if old_text in para.text:
                        for run in para.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                                changed = True
                else:
                    if old_text.lower() in para.text.lower():
                        for run in para.runs:
                            if old_text.lower() in run.text.lower():
                                run.text = run.text.replace(old_text, new_text)
                                changed = True
            if not changed:
                return "未找到匹配文本，未做修改"
        elif action == "add_paragraph":
            if not text:
                return "错误：add_paragraph 需要提供 text"
            doc.add_paragraph(text)
            changed = True
        doc.save(path)
        return f"✅ 修改成功：{action} 操作已完成"
    except Exception as e:
        return f"❌ 修改失败：{str(e)}"


def edit_docx_by_sample(path, sample, replacement):
    if not path or not sample or not replacement:
        return "错误：需要提供 path、sample 和 replacement"
    if not os.path.exists(path):
        return f"错误：文件不存在: {path}"
    if not path.lower().endswith('.docx'):
        return f"错误：不是 DOCX 文件: {path}"
    try:
        doc = Document(path)
        modified = False
        for para in doc.paragraphs:
            runs_text = "".join(run.text for run in para.runs)
            if sample not in runs_text:
                continue
            start_idx = runs_text.index(sample)
            end_idx = start_idx + len(sample)
            offset = 0
            target_indices = []
            for i, run in enumerate(para.runs):
                run_start = offset
                run_end = offset + len(run.text)
                if run_start < end_idx and run_end > start_idx:
                    target_indices.append(i)
                offset = run_end
            if not target_indices:
                continue
            para.runs[target_indices[0]].text = replacement
            for idx in target_indices[1:]:
                para.runs[idx].text = ""
            modified = True
            break
        if modified:
            doc.save(path)
            return f"✅ 跨run替换成功：已将「{sample}」替换为指定文本"
        else:
            return "未找到匹配文本，未做修改"
    except Exception as e:
        return f"❌ 替换失败：{str(e)}"


def batch_read(directory, recursive=False):
    if not os.path.isdir(directory):
        return f"错误：目录不存在: {directory}"
    if recursive:
        docx_files = []
        for root, dirs, files in os.walk(directory):
            for f in files:
                if f.lower().endswith('.docx') and not f.startswith('~$'):
                    docx_files.append(os.path.join(root, f))
    else:
        docx_files = [os.path.join(directory, f) for f in os.listdir(directory) if f.lower().endswith('.docx') and not f.startswith('~$')]
    docx_files.sort()
    if not docx_files:
        return f"在目录中未找到 DOCX 文件: {directory}"
    results = []
    success = 0
    failed = 0
    for docx_path in docx_files:
        try:
            doc = Document(docx_path)
            para_count = len([p for p in doc.paragraphs if p.text.strip()])
            results.append(f"✅ {os.path.basename(docx_path)} ({para_count} 个非空段落)")
            success += 1
        except Exception as e:
            results.append(f"❌ {os.path.basename(docx_path)}: {str(e)}")
            failed += 1
    return f"批量读取完成：成功 {success} 个，失败 {failed} 个，共 {len(docx_files)} 个文件\n" + '\n'.join(results)


def read_pdf(path):
    if not path:
        return "错误：请提供 PDF 文件路径"
    if not os.path.exists(path):
        return f"错误：文件不存在: {path}"
    if not path.lower().endswith('.pdf'):
        return f"错误：不是 PDF 文件: {path}"
    try:
        import fitz
        with fitz.open(path) as doc:
            lines = [f"📄 文件: {os.path.basename(path)}", f"📄 共 {len(doc)} 页", '=' * 50]
            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                lines.append(f"\n【第 {i} 页】")
                lines.append(text if text else "[空白页或无文本]")
            return '\n'.join(lines)
    except ImportError:
        return "错误：需要安装 PyMuPDF，请运行 pip install pymupdf"
    except Exception as e:
        return f"❌ PDF提取失败: {str(e)}"
